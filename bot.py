import os
import re
import logging
import tempfile
import httpx
import trafilatura
from dotenv import load_dotenv
from supabase import create_client
from telegram import Update
from telegram.ext import Application, MessageHandler, CommandHandler, filters, ContextTypes

load_dotenv()

logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logging.getLogger("httpx").setLevel(logging.WARNING)
log = logging.getLogger("hermes")

TOKEN = os.environ["TELEGRAM_BOT_TOKEN"]
ALLOWED = os.getenv("ALLOWED_USER_ID", "").strip()
NVIDIA_KEY = os.environ["NVIDIA_API_KEY"]
OPENROUTER_KEY = os.environ["OPENROUTER_API_KEY"]

EMBED_URL = "https://integrate.api.nvidia.com/v1/embeddings"
EMBED_MODEL = "nvidia/nv-embedqa-e5-v5"
CHAT_URL = "https://openrouter.ai/api/v1/chat/completions"
CHAT_MODEL = "anthropic/claude-sonnet-5"

CHUNK_WORDS = 250
CHUNK_OVERLAP = 40
BATCH_SIZE = 16

URL_RE = re.compile(r"(?:https?://|www\.)[^\s<>\"']+", re.IGNORECASE)

SYSTEM_PROMPT = (
    "You are Hermes, a personal memory assistant. Answer using the retrieved "
    "memories provided. Be concise and direct. If the memories don't contain "
    "the answer, say so plainly rather than guessing. Cite the source name "
    "and date where useful."
)

supabase = create_client(
    os.environ["SUPABASE_URL"],
    os.environ["SUPABASE_SERVICE_KEY"],
)


# ---------- sanitising ----------

def clean_text(text: str) -> str:
    """Strip NUL bytes and control chars Postgres refuses to store."""
    if not text:
        return ""
    text = text.replace("\x00", "")
    text = "".join(ch for ch in text if ch in "\n\t" or ord(ch) >= 32)
    return text


def normalise_url(url: str) -> str:
    url = url.rstrip(".,);:'\"")
    if not url.lower().startswith(("http://", "https://")):
        url = "https://" + url
    return url


# ---------- text extraction ----------

def extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n\n".join((p.extract_text() or "") for p in reader.pages)


def extract_docx(path: str) -> str:
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        bits = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                bits.append(f"(notes) {notes}")
        if bits:
            parts.append(f"Slide {i}: " + "\n".join(bits))
    return "\n\n".join(parts)


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".txt": lambda p: open(p, encoding="utf-8", errors="ignore").read(),
    ".md": lambda p: open(p, encoding="utf-8", errors="ignore").read(),
}


# ---------- chunking ----------

def chunk_text(text: str) -> list[str]:
    text = clean_text(text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    words = text.split()
    if not words:
        return []
    if len(words) <= CHUNK_WORDS:
        return [" ".join(words)]

    chunks = []
    step = CHUNK_WORDS - CHUNK_OVERLAP
    for i in range(0, len(words), step):
        piece = " ".join(words[i:i + CHUNK_WORDS])
        if len(piece.split()) > 20:
            chunks.append(piece)
    return chunks


# ---------- embedding ----------

async def embed_batch(texts: list[str], input_type: str) -> list[list[float]]:
    out = []
    async with httpx.AsyncClient(timeout=60) as client:
        for i in range(0, len(texts), BATCH_SIZE):
            batch = texts[i:i + BATCH_SIZE]
            r = await client.post(
                EMBED_URL,
                headers={"Authorization": f"Bearer {NVIDIA_KEY}"},
                json={
                    "input": batch,
                    "model": EMBED_MODEL,
                    "input_type": input_type,
                    "truncate": "END",
                },
            )
            r.raise_for_status()
            data = sorted(r.json()["data"], key=lambda d: d["index"])
            out.extend(d["embedding"] for d in data)
    return out


async def embed(text: str, input_type: str) -> list[float]:
    return (await embed_batch([text], input_type))[0]


# ---------- storing ----------

async def store_chunks(chunks: list[str], source: str, meta: dict) -> int:
    if not chunks:
        return 0
    vectors = await embed_batch(chunks, "passage")
    rows = [
        {
            "content": c,
            "embedding": v,
            "source": source,
            "metadata": {**meta, "chunk": i + 1, "chunks": len(chunks)},
        }
        for i, (c, v) in enumerate(zip(chunks, vectors))
    ]
    for i in range(0, len(rows), 50):
        supabase.table("memories").insert(rows[i:i + 50]).execute()
    return len(rows)


# ---------- retrieval ----------

async def search(query: str, limit: int = 8) -> list[dict]:
    vec = await embed(query, "query")
    res = supabase.rpc("match_memories", {
        "query_embedding": vec,
        "match_threshold": 0.25,
        "match_count": limit,
    }).execute()
    return res.data or []


async def reason(question: str, memories: list[dict]) -> str:
    if memories:
        parts = []
        for m in memories:
            label = m["metadata"].get("title") or m["metadata"].get("file_name") or m["source"]
            parts.append(f"[{m['created_at'][:10]} · {label}] {m['content']}")
        context = "\n\n".join(parts)
    else:
        context = "(no relevant memories found)"

    async with httpx.AsyncClient(timeout=90) as client:
        r = await client.post(
            CHAT_URL,
            headers={
                "Authorization": f"Bearer {OPENROUTER_KEY}",
                "HTTP-Referer": "https://circularsmart.com",
                "X-Title": "Hermes",
            },
            json={
                "model": CHAT_MODEL,
                "messages": [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content":
                        f"Retrieved memories:\n\n{context}\n\n---\n\nQuestion: {question}"},
                ],
                "max_tokens": 1200,
            },
        )
        r.raise_for_status()
        return r.json()["choices"][0]["message"]["content"]


# ---------- handlers ----------

def authorised(update: Update) -> bool:
    if not ALLOWED:
        return True
    return str(update.effective_user.id) == ALLOWED


async def whoami(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(f"Your Telegram user ID is: {update.effective_user.id}")


async def recall(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not authorised(update):
        return
    query = " ".join(context.args).strip()
    if not query:
        await update.message.reply_text("Usage: /recall search terms")
        return
    try:
        rows = await search(query, limit=5)
        if not rows:
            await update.message.reply_text("Nothing relevant found.")
            return
        lines = []
        for r in rows:
            label = r["metadata"].get("title") or r["metadata"].get("file_name") or r["source"]
            lines.append(f"[{r['created_at'][:10]} · {label} · {round(r['similarity'], 2)}]\n{r['content'][:400]}")
        await update.message.reply_text("\n\n".join(lines)[:4000])
    except Exception as e:
        log.exception("Recall failed")
        await update.message.reply_text(f"Recall failed: {e}")


async def ask(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not authorised(update):
        return
    question = " ".join(context.args).strip()
    if not question:
        await update.message.reply_text("Usage: /ask your question")
        return
    await update.message.chat.send_action("typing")
    try:
        rows = await search(question)
        answer = await reason(question, rows)
        await update.message.reply_text(answer[:4000])
    except Exception as e:
        log.exception("Ask failed")
        await update.message.reply_text(f"Ask failed: {e}")


async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not authorised(update):
        return

    doc = update.message.document
    name = doc.file_name or "unnamed"
    ext = os.path.splitext(name)[1].lower()

    if ext not in EXTRACTORS:
        await update.message.reply_text(
            f"Can't read {ext or 'that'}. Supported: pdf, docx, pptx, txt, md"
        )
        return

    await update.message.reply_text(f"Reading {name}...")
    await update.message.chat.send_action("typing")

    try:
        tg_file = await context.bot.get_file(doc.file_id)
        with tempfile.NamedTemporaryFile(suffix=ext, delete=True) as tmp:
            await tg_file.download_to_drive(tmp.name)
            text = EXTRACTORS[ext](tmp.name)

        chunks = chunk_text(text)
        if not chunks:
            await update.message.reply_text("No readable text found — is it a scanned image?")
            return

        n = await store_chunks(chunks, ext.lstrip("."), {
            "file_name": name,
            "title": os.path.splitext(name)[0],
            "user_id": update.effective_user.id,
        })
        await update.message.reply_text(f"Stored {name} as {n} chunks.")

    except Exception as e:
        log.exception("Document ingest failed")
        await update.message.reply_text(f"Failed: {e}")


async def ingest_url(update: Update, url: str):
    await update.message.reply_text(f"Fetching {url[:60]}...")
    await update.message.chat.send_action("typing")

    try:
        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            await update.message.reply_text(f"Couldn't fetch {url[:60]}")
            return

        text = trafilatura.extract(downloaded, include_comments=False, include_tables=True)
        meta = trafilatura.extract_metadata(downloaded)
        title = (meta.title if meta else None) or url

        if not text:
            await update.message.reply_text("No readable content extracted.")
            return

        chunks = chunk_text(text)
        n = await store_chunks(chunks, "web", {
            "url": url,
            "title": clean_text(title)[:200],
            "user_id": update.effective_user.id,
        })
        await update.message.reply_text(f"Stored \"{title[:80]}\" as {n} chunks.")

    except Exception as e:
        log.exception("URL ingest failed")
        await update.message.reply_text(f"Failed: {e}")


async def save_note(update: Update, text: str):
    vec = await embed(text, "passage")
    supabase.table("memories").insert({
        "content": text,
        "embedding": vec,
        "source": "telegram",
        "metadata": {
            "user_id": update.effective_user.id,
            "chat_id": update.effective_chat.id,
            "message_id": update.message.message_id,
        },
    }).execute()


async def capture(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not authorised(update):
        log.warning("Rejected message from %s", update.effective_user.id)
        return

    text = clean_text(update.message.text)
    urls = [normalise_url(u) for u in URL_RE.findall(text)]

    try:
        if urls:
            # fetch every link in the message
            for url in urls[:3]:
                await ingest_url(update, url)

            # if there's real commentary alongside the link, keep that too
            remainder = URL_RE.sub("", text).strip()
            if len(remainder.split()) >= 5:
                await save_note(update, text)
                await update.message.reply_text("Note saved alongside.")
            return

        await save_note(update, text)
        await update.message.reply_text("Saved.")

    except Exception as e:
        log.exception("Capture failed")
        await update.message.reply_text(f"Failed to save: {e}")


def main():
    app = Application.builder().token(TOKEN).build()
    app.add_handler(CommandHandler("whoami", whoami))
    app.add_handler(CommandHandler("recall", recall))
    app.add_handler(CommandHandler("ask", ask))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, capture))
    log.info("Hermes starting (polling)")
    app.run_polling()


if __name__ == "__main__":
    main()
