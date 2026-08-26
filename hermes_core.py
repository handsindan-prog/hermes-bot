"""
hermes_core — the shared half of Hermes.

Config, sanitising, extraction, chunking, embedding, dedupe-aware storage,
retrieval, reasoning and run logging. The Telegram bot and every scheduled
agent import from here; nothing in this module knows Telegram exists.

Extracted from bot.py so that agent number two does not begin its life by
copy-pasting chunk/embed/store for the third time.
"""

from __future__ import annotations

import hashlib
import logging
import os
import re
from dataclasses import dataclass, field
from datetime import datetime, timedelta, timezone
from pathlib import Path

import httpx
from dotenv import load_dotenv
from supabase import create_client

# .env sits beside this file. Resolve it explicitly rather than relying on
# find_dotenv walking up from the cwd, which breaks for anything launched by
# a systemd timer from a different working directory.
load_dotenv(Path(__file__).resolve().parent / ".env")

log = logging.getLogger("hermes")

# ---------- config ----------

NVIDIA_KEY = os.environ["NVIDIA_API_KEY"]
OPENROUTER_KEY = os.environ["OPENROUTER_API_KEY"]

EMBED_URL = "https://integrate.api.nvidia.com/v1/embeddings"
# nv-embedqa-e5-v5 reached end of life on 2026-08-25T09:00Z and now returns 410.
# nemotron-3-embed-1b is the only embedding model this account can actually call
# — the others in /v1/models 404 — and it is asymmetric like the old one.
#
# 2048 dims, which the original build rejected because ivfflat and hnsw cap at
# 2000. That no longer applies: the index was dropped in migration 002 after it
# turned out to be suppressing search, and exact search has no dimension limit.
#
# Changing this model invalidates every stored vector. They live in a different
# space and are not comparable, so the corpus must be re-embedded — see
# agents/reembed.py.
# Third embedding model in 24 hours. nv-embedqa-e5-v5 was retired (410),
# nemotron-3-embed-1b then began hanging — 3/3 ReadTimeouts at 30s while NIM
# chat stayed healthy, so it is that model's endpoint, not the account.
# llama-nemotron-embed-vl-1b-v2 also emits 2048 dims, so no schema change.
#
# NIM has proven unreliable for embeddings specifically. If this recurs, move
# the embedding provider off NIM rather than picking a fourth model.
EMBED_MODEL = os.getenv("HERMES_EMBED_MODEL", "nvidia/llama-nemotron-embed-vl-1b-v2")
EMBED_DIMS = 2048

# Two tiers on purpose. An agent sweeping thirty pages is ~90% mechanical
# work — classify, extract, first-pass summarise — and paying Claude rates
# for it is the expensive way to do a cheap job. "smart" is for the single
# final judgement; "fast" is for the volume underneath it.
#
# The fast default is measured, not guessed: on an identical classification
# task nemotron-3-nano-30b-a3b spent 53 completion tokens where
# nemotron-nano-9b-v2 spent 203 and nemotron-super-49b-v1.5 spent 271, all
# three returning the same correct answer. Note that
# llama-3.1-nemotron-70b-instruct appears in /v1/models but 404s on this
# account — check a model actually answers before adopting it.
SMART_URL = "https://openrouter.ai/api/v1/chat/completions"
SMART_MODEL = os.getenv("HERMES_SMART_MODEL", "anthropic/claude-sonnet-5")
FAST_URL = "https://integrate.api.nvidia.com/v1/chat/completions"
FAST_MODEL = os.getenv("HERMES_FAST_MODEL", "nvidia/nemotron-3-nano-30b-a3b")

CHUNK_WORDS = 250          # sized for the embedding model's 512-token limit
CHUNK_OVERLAP = 40
BATCH_SIZE = 16
MIN_CHUNK_WORDS = 20
INSERT_BATCH = 50
MATCH_THRESHOLD = 0.25     # real matches cluster ~0.35; the useful band is narrow

URL_RE = re.compile(r"(?:https?://|www\.)[^\s<>\"']+", re.IGNORECASE)

SYSTEM_PROMPT = (
    "You are Hermes, a personal memory assistant. Answer using the retrieved "
    "memories provided. Be concise and direct. If the memories don't contain "
    "the answer, say so plainly rather than guessing. Cite the source name "
    "and date where useful."
)

# Which product/entity this process writes on behalf of. One Hermes serves a
# portfolio; rows without a workspace would be unattributable later.
WORKSPACE = os.getenv("HERMES_WORKSPACE", "circularsmart")

supabase = create_client(
    os.environ["SUPABASE_URL"],
    os.environ["SUPABASE_SERVICE_KEY"],
)


# ---------- sanitising ----------

def clean_text(text: str) -> str:
    """Strip NUL bytes and control chars Postgres refuses to store."""
    if not text:
        return ""
    text = text.replace("\x00", "")
    return "".join(ch for ch in text if ch in "\n\t" or ord(ch) >= 32)


def normalise_url(url: str) -> str:
    url = url.rstrip(".,);:'\"")
    if not url.lower().startswith(("http://", "https://")):
        url = "https://" + url
    return url


def sha256(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", "replace")).hexdigest()


# ---------- text extraction ----------

def extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n\n".join((p.extract_text() or "") for p in reader.pages)


def extract_docx(path: str) -> str:
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


# A single sheet can hold tens of thousands of rows. Past a few hundred it is
# a database, not a document, and indexing it whole would swamp the corpus with
# one file. Truncation is reported in the text rather than done silently.
MAX_SHEET_ROWS = 500


def extract_xlsx(path: str) -> str:
    """Spreadsheets as readable text — one block per sheet, rows pipe-delimited.

    data_only=True reads the values Excel caches alongside each formula. The
    formula string "=SUM(B2:B40)" is noise in a semantic index; the number it
    produced is not. A workbook written by a tool that never cached its values
    extracts as empty, which library_sync reports as "no extractable text"
    rather than storing a blank document.
    """
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True, read_only=True)
    try:
        blocks = []
        for ws in wb.worksheets:
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= MAX_SHEET_ROWS:
                    rows.append(f"(… further rows beyond {MAX_SHEET_ROWS} not indexed)")
                    break
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                # The sheet name is often the only thing naming what the numbers
                # are, so it goes in the text rather than the metadata.
                blocks.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
        return "\n\n".join(blocks)
    finally:
        wb.close()


def extract_csv(path: str) -> str:
    import csv
    with open(path, newline="", encoding="utf-8", errors="ignore") as fh:
        sample = fh.read(8192)
        fh.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except csv.Error:
            dialect = csv.excel
        rows = []
        for i, row in enumerate(csv.reader(fh, dialect)):
            if i >= MAX_SHEET_ROWS:
                rows.append(f"(… further rows beyond {MAX_SHEET_ROWS} not indexed)")
                break
            cells = [c.strip() for c in row if c and c.strip()]
            if cells:
                rows.append(" | ".join(cells))
    return "\n".join(rows)


def extract_html(path: str) -> str:
    """Saved web pages — trafilatura strips the chrome the same way it does for
    a live fetch, so a saved article reads like a fetched one."""
    import trafilatura
    raw = open(path, encoding="utf-8", errors="ignore").read()
    return trafilatura.extract(raw, include_comments=False, include_tables=True) or ""


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        bits = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                bits.append(f"(notes) {notes}")
        if bits:
            parts.append(f"Slide {i}: " + "\n".join(bits))
    return "\n\n".join(parts)


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".htm": extract_html,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".csv": extract_csv,
    ".txt": lambda p: open(p, encoding="utf-8", errors="ignore").read(),
    ".md": lambda p: open(p, encoding="utf-8", errors="ignore").read(),
}


# ---------- chunking ----------

def chunk_text(text: str) -> list[str]:
    text = clean_text(text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    words = text.split()
    if not words:
        return []
    if len(words) <= CHUNK_WORDS:
        return [" ".join(words)]

    chunks = []
    step = CHUNK_WORDS - CHUNK_OVERLAP
    for i in range(0, len(words), step):
        piece = " ".join(words[i:i + CHUNK_WORDS])
        if len(piece.split()) > MIN_CHUNK_WORDS:
            chunks.append(piece)
    return chunks


# ---------- embedding ----------

async def embed_batch(texts: list[str], input_type: str) -> list[list[float]]:
    """
    input_type matters and fails silently if wrong: 'passage' when storing,
    'query' when searching. Mismatched types wreck recall without erroring.
    """
    out: list[list[float]] = []
    async with httpx.AsyncClient(timeout=60) as client:
        for i in range(0, len(texts), BATCH_SIZE):
            batch = texts[i:i + BATCH_SIZE]
            r = await client.post(
                EMBED_URL,
                headers={"Authorization": f"Bearer {NVIDIA_KEY}"},
                json={
                    "input": batch,
                    "model": EMBED_MODEL,
                    "input_type": input_type,
                    "truncate": "END",
                },
            )
            r.raise_for_status()
            data = sorted(r.json()["data"], key=lambda d: d["index"])
            out.extend(d["embedding"] for d in data)
    return out


async def embed(text: str, input_type: str) -> list[float]:
    return (await embed_batch([text], input_type))[0]


# ---------- storing ----------

@dataclass
class IngestResult:
    status: str          # stored | updated | unchanged | empty
    written: int = 0
    skipped: int = 0
    doc_key: str = ""
    workspace: str = ""

    def __str__(self) -> str:
        if self.status == "unchanged":
            return f"{self.doc_key}: unchanged, skipped {self.skipped} chunks"
        if self.status == "empty":
            return f"{self.doc_key}: no readable text"
        return f"{self.doc_key}: {self.status} {self.written} chunks"


def _dedupe_chunks(chunks: list[str]) -> list[tuple[str, str]]:
    """Drop chunks identical to an earlier one in the same document, before
    they cost an embedding call. Returns (chunk, content_hash) pairs."""
    seen: set[str] = set()
    out = []
    for c in chunks:
        h = sha256(c)
        if h in seen:
            continue
        seen.add(h)
        out.append((c, h))
    return out


async def store_document(
    doc_key: str,
    text: str,
    source: str,
    meta: dict | None = None,
    workspace: str | None = None,
) -> IngestResult:
    """
    Idempotent document ingestion.

    Re-ingesting an unchanged document is free and writes nothing. Re-ingesting
    a changed one replaces its chunks rather than adding a second copy. This is
    what makes a scheduled Drive or web sync safe to run repeatedly — without
    it, every run doubles the corpus and duplicate chunks crowd genuine answers
    out of match_memories.
    """
    meta = dict(meta or {})
    ws = workspace or WORKSPACE
    chunks = chunk_text(text)
    if not chunks:
        return IngestResult("empty", doc_key=doc_key, workspace=ws)

    doc_hash = sha256(clean_text(text))

    # Document identity is (workspace, doc_key) — the same competitor page can
    # legitimately sit in two products' corpora without one clobbering the other.
    existing = (
        supabase.table("memories")
        .select("doc_hash")
        .eq("workspace", ws)
        .eq("doc_key", doc_key)
        .limit(1)
        .execute()
        .data
    )

    if existing and existing[0].get("doc_hash") == doc_hash:
        return IngestResult("unchanged", skipped=len(chunks), doc_key=doc_key, workspace=ws)

    status = "stored"
    if existing:
        # Known document, new content: replace rather than accumulate.
        supabase.table("memories").delete().eq("workspace", ws).eq("doc_key", doc_key).execute()
        status = "updated"

    pairs = _dedupe_chunks(chunks)
    vectors = await embed_batch([c for c, _ in pairs], "passage")

    rows = [
        {
            "content": c,
            "embedding": v,
            "source": source,
            "workspace": ws,
            "doc_key": doc_key,
            "doc_hash": doc_hash,
            "content_hash": h,
            "metadata": {**meta, "chunk": i + 1, "chunks": len(pairs)},
        }
        for i, ((c, h), v) in enumerate(zip(pairs, vectors))
    ]

    for i in range(0, len(rows), INSERT_BATCH):
        (
            supabase.table("memories")
            .upsert(
                rows[i:i + INSERT_BATCH],
                on_conflict="workspace,doc_key,content_hash",
                ignore_duplicates=True,
            )
            .execute()
        )

    return IngestResult(
        status,
        written=len(rows),
        skipped=len(chunks) - len(pairs),
        doc_key=doc_key,
        workspace=ws,
    )


async def store_note(content: str, meta: dict | None = None, source: str = "telegram",
                     workspace: str | None = None) -> int:
    """A loose note has no document identity; the DB trigger assigns it
    doc_key = 'note:<id>' once the row has an id."""
    vec = await embed(content, "passage")
    supabase.table("memories").insert({
        "content": content,
        "embedding": vec,
        "source": source,
        "workspace": workspace or WORKSPACE,
        "metadata": meta or {},
    }).execute()
    return 1


# ---------- retrieval ----------

async def search(
    query: str,
    limit: int = 8,
    threshold: float = MATCH_THRESHOLD,
    workspace: str | None = None,
    source: str | None = None,
) -> list[dict]:
    """workspace=None searches the whole corpus; pass one to scope to a product."""
    vec = await embed(query, "query")
    res = supabase.rpc("match_memories", {
        "query_embedding": vec,
        "match_threshold": threshold,
        "match_count": limit,
        "filter_source": source,
        "filter_workspace": workspace,
    }).execute()
    return res.data or []


def label_of(row: dict) -> str:
    meta = row.get("metadata") or {}
    return meta.get("title") or meta.get("file_name") or row.get("source", "?")


def as_context(memories: list[dict]) -> str:
    if not memories:
        return "(no relevant memories found)"
    return "\n\n".join(
        f"[{m['created_at'][:10]} · {label_of(m)}] {m['content']}"
        for m in memories
    )


# ---------- reasoning ----------

class ChatTruncated(RuntimeError):
    """The model hit max_tokens before finishing its answer.

    Worth its own exception because of how these models fail. A starved
    Nemotron call does not return null — it returns its own chain-of-thought
    in the content field, mid-sentence:

        max_tokens=8 → content='We need to classify as one word,'

    That reads like prose and would be written straight into the corpus as an
    answer. Any completion that stopped on 'length' is therefore treated as no
    answer at all, not as a short one.
    """


@dataclass
class ChatResult:
    text: str
    input_tokens: int = 0
    output_tokens: int = 0
    cost_usd: float = 0.0
    model: str = ""
    reasoning: str = ""
    truncated: bool = False


async def chat(
    messages: list[dict],
    tier: str = "smart",
    max_tokens: int = 1200,
    timeout: int = 90,
    strict: bool | None = None,
) -> ChatResult:
    """
    tier='smart' → Claude via OpenRouter, for judgement and final prose.
    tier='fast'  → Nemotron via NVIDIA NIM, for high-volume mechanical steps.

    Both endpoints are OpenAI-shaped. OpenRouter is asked to report actual
    spend so agent_runs.cost_usd is measured rather than estimated; NIM does
    not report cost, so fast-tier calls log tokens with cost 0.

    Every Nemotron model on NIM is a reasoning model: it emits
    reasoning_content *and* content, and the reasoning is billed against
    max_tokens. Starve it and content comes back null with
    finish_reason='length' — so do not set max_tokens tight on the fast tier
    just because the answer is one word. That case raises ChatTruncated
    rather than returning an empty string, because a silent "" is the kind of
    thing an agent will happily write into the corpus.
    """
    if tier == "fast":
        url, model, key = FAST_URL, FAST_MODEL, NVIDIA_KEY
        payload = {"model": model, "messages": messages, "max_tokens": max_tokens}
        headers = {"Authorization": f"Bearer {key}"}
    else:
        url, model, key = SMART_URL, SMART_MODEL, OPENROUTER_KEY
        payload = {
            "model": model,
            "messages": messages,
            "max_tokens": max_tokens,
            "usage": {"include": True},
        }
        headers = {
            "Authorization": f"Bearer {key}",
            "HTTP-Referer": "https://circularsmart.com",
            "X-Title": "Hermes",
        }

    async with httpx.AsyncClient(timeout=timeout) as client:
        r = await client.post(url, headers=headers, json=payload)
        r.raise_for_status()
        body = r.json()

    usage = body.get("usage") or {}
    choice = body["choices"][0]
    msg = choice.get("message") or {}

    text = (msg.get("content") or "").strip()
    reasoning = (msg.get("reasoning_content") or msg.get("reasoning") or "").strip()
    truncated = choice.get("finish_reason") == "length"

    if strict is None:
        strict = tier == "fast"

    if truncated and strict:
        raise ChatTruncated(
            f"{model} stopped at the {max_tokens}-token limit after "
            f"{usage.get('completion_tokens', '?')} tokens. What it returned is an "
            f"unfinished thought, not an answer: {text[:80]!r}. Raise max_tokens."
        )

    if not text:
        # Some models answer entirely inside the reasoning channel.
        text = reasoning

    return ChatResult(
        text=text,
        input_tokens=usage.get("prompt_tokens", 0) or 0,
        output_tokens=usage.get("completion_tokens", 0) or 0,
        cost_usd=float(usage.get("cost", 0) or 0),
        model=model,
        reasoning=reasoning,
        truncated=truncated,
    )


async def reason(question: str, memories: list[dict], tier: str = "smart") -> ChatResult:
    return await chat(
        [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content":
                f"Retrieved memories:\n\n{as_context(memories)}\n\n---\n\nQuestion: {question}"},
        ],
        tier=tier,
    )


# ---------- findings ----------
#
# Deliberately not "claim": Claim is one of the six Vantage modules, and the
# research agent's own report used the word in both senses at once. A finding
# is what it is — something discovered, carrying evidence.

CONFIDENCE = ("verified", "high", "medium", "low", "unverified")


@dataclass
class Finding:
    """One thing an agent believes, and why.

    `evidence` is the verbatim text the finding rests on — not a paraphrase.
    Without it a reviewer cannot check the finding without redoing the research,
    which defeats the point of the review.

    A finding whose source_kind is 'model' is forced to 'unverified' by a
    database constraint no matter what confidence is passed, because a model's
    certainty about an entity is not evidence about that entity.
    """

    target: str                     # brain file this answers into
    statement: str                  # what is asserted, in one sentence
    marker: str = ""                # which [NEEDS RESEARCH] it addresses
    evidence: str = ""              # verbatim quote
    source_url: str = ""
    source_kind: str = "corpus"     # corpus | web | apify | model
    confidence: str = "unverified"
    fetched_at: str | None = None

    def __post_init__(self):
        if self.confidence not in CONFIDENCE:
            raise ValueError(f"confidence must be one of {CONFIDENCE}, got {self.confidence!r}")
        if self.source_kind == "model":
            self.confidence = "unverified"


def record_finding(f: Finding, run_id: int | None = None, workspace: str | None = None) -> int:
    row = supabase.table("agent_findings").insert({
        "run_id": run_id,
        "workspace": workspace or WORKSPACE,
        "target": f.target,
        "marker": f.marker or None,
        "statement": f.statement,
        "evidence": f.evidence or None,
        "source_url": f.source_url or None,
        "source_kind": f.source_kind,
        "confidence": f.confidence,
        "fetched_at": f.fetched_at,
    }).execute().data
    return row[0]["id"] if row else 0


def open_findings(workspace: str | None = None, target: str | None = None) -> list[dict]:
    q = (supabase.table("agent_findings").select("*")
         .eq("workspace", workspace or WORKSPACE)
         .eq("status", "proposed")
         .order("created_at", desc=True))
    if target:
        q = q.eq("target", target)
    return q.execute().data or []


# ---------- run logging ----------

class SpendCapExceeded(RuntimeError):
    pass


@dataclass
class AgentRun:
    """
    Context manager wrapping one agent execution. Writes a row to agent_runs on
    entry and closes it on exit, whatever happens — so a crashed agent leaves
    'failed' with a traceback rather than silence.

    Read-only by default is a convention, not a mechanism; the spend cap is the
    mechanism. check_cap() consults spend across the whole window, not just this
    run, so a wedged agent retrying every five minutes still halts.

        async with AgentRun("competitor-watch", spend_cap_usd=0.50) as run:
            run.fetched(3)
            res = await hermes_core.chat(msgs); run.charge(res)
            run.check_cap()
    """

    agent: str
    workspace: str = ""
    spend_cap_usd: float | None = None
    cap_window_hours: int = 24
    detail: dict = field(default_factory=dict)

    run_id: int | None = None
    sources_fetched: int = 0
    rows_written: int = 0
    rows_skipped: int = 0
    input_tokens: int = 0
    output_tokens: int = 0
    cost_usd: float = 0.0

    # -- accounting --

    def fetched(self, n: int = 1) -> None:
        self.sources_fetched += n

    def wrote(self, result: IngestResult | int) -> None:
        if isinstance(result, IngestResult):
            self.rows_written += result.written
            self.rows_skipped += result.skipped
        else:
            self.rows_written += result

    def charge(self, result: ChatResult) -> None:
        self.input_tokens += result.input_tokens
        self.output_tokens += result.output_tokens
        self.cost_usd += result.cost_usd

    def finding(self, f: "Finding") -> int:
        """Record a finding for human review. Never writes to memories — a
        finding is not a fact until somebody accepts it."""
        return record_finding(f, run_id=self.run_id, workspace=self.ws)

    @property
    def ws(self) -> str:
        return self.workspace or WORKSPACE

    def window_spend(self) -> float:
        """Spend across the window for this workspace — not just this run, so an
        agent wedged in a retry loop still trips its own cap."""
        since = (datetime.now(timezone.utc)
                 - timedelta(hours=self.cap_window_hours)).isoformat()
        res = supabase.rpc("agent_spend_since", {
            "since": since,
            "filter_workspace": self.ws,
        }).execute()
        return float(res.data or 0)

    def check_cap(self) -> None:
        """Raise if this run has pushed spend over the cap. Call between steps."""
        if self.spend_cap_usd is None:
            return
        spent = self.window_spend() + self.cost_usd
        if spent > self.spend_cap_usd:
            raise SpendCapExceeded(
                f"{self.agent}: ${spent:.4f} spent in the last "
                f"{self.cap_window_hours}h exceeds cap ${self.spend_cap_usd:.2f}"
            )

    # -- lifecycle --

    async def __aenter__(self) -> "AgentRun":
        row = supabase.table("agent_runs").insert({
            "agent": self.agent,
            "workspace": self.ws,
            "status": "running",
            "detail": self.detail,
        }).execute().data
        self.run_id = row[0]["id"] if row else None
        log.info("%s: run %s started", self.agent, self.run_id)
        return self

    async def __aexit__(self, exc_type, exc, tb) -> bool:
        if exc_type is None:
            status, error = "ok", None
        elif exc_type is SpendCapExceeded:
            status, error = "halted", str(exc)
        else:
            status, error = "failed", f"{exc_type.__name__}: {exc}"

        if self.run_id is not None:
            supabase.table("agent_runs").update({
                "status": status,
                "finished_at": datetime.now(timezone.utc).isoformat(),
                "sources_fetched": self.sources_fetched,
                "rows_written": self.rows_written,
                "rows_skipped": self.rows_skipped,
                "input_tokens": self.input_tokens,
                "output_tokens": self.output_tokens,
                "cost_usd": round(self.cost_usd, 6),
                "detail": self.detail,
                "error": error,
            }).eq("id", self.run_id).execute()

        log.info(
            "%s: run %s %s — %s fetched, %s written, $%.4f",
            self.agent, self.run_id, status,
            self.sources_fetched, self.rows_written, self.cost_usd,
        )

        # A spend cap stopping the run is a designed outcome, not a crash.
        return exc_type is SpendCapExceeded
